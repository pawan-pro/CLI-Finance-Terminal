import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from datetime import datetime, timedelta
import os
import time
from polygon import RESTClient
import pytz
import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
import MetaTrader5 as mt5
from mplfinance.original_flavor import candlestick_ohlc
from mpl_toolkits.axes_grid1 import make_axes_locatable
import re
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Polygon API key
API_KEY = "f8Oowzm4TdXsQjxOOlH7C1_2xB1dxNSB"

# Initialize Polygon RESTClient
client = RESTClient(api_key=API_KEY)

# Retry strategy to handle 429 errors
retry_strategy = Retry(
    total=2,
    backoff_factor=2,
    status_forcelist=[429, 500, 502, 503, 504],
    allowed_methods=["GET"]
)
adapter = HTTPAdapter(max_retries=retry_strategy)
http = requests.Session()
http.mount("https://", adapter)

# Initialize MT5 connection
if not mt5.initialize():
    print("MT5 initialization failed")
    mt5.shutdown()

def estimate_typical_spread(symbol):
    major_fx = {
        "EURUSD.sd": 0.00010, "USDJPY.sd": 0.00010, "GBPUSD.sd": 0.00015,
        "USDCHF.sd": 0.00015, "AUDUSD.sd": 0.00015, "USDCAD.sd": 0.00015
    }
    minor_fx = {
        "EURGBP.sd": 0.00020, "EURJPY.sd": 0.00020, "GBPJPY.sd": 0.00025,
        "EURCHF.sd": 0.00025, "AUDJPY.sd": 0.00020, "NZDUSD.sd": 0.00020
    }
    
    if symbol in major_fx:
        return major_fx[symbol]
    elif symbol in minor_fx:
        return minor_fx[symbol]
    else:
        return 0.00030  # Default to 3 pips

def get_current_spread(symbol):
    symbol_info = mt5.symbol_info(symbol)
    if symbol_info is None:
        print(f"Failed to get symbol info for {symbol}")
        return None
    
    spread = symbol_info.ask - symbol_info.bid
    if spread == 0:
        typical_spread = estimate_typical_spread(symbol)
        print(f"MT5 returned 0 spread for {symbol}. Using estimated typical spread: {typical_spread}")
        return typical_spread
    
    return spread

def fetch_data(symbol, start_time, end_time):
    if symbol == "USOILRoll":
        # Convert to UTC for MT5
        start_time_utc = start_time.astimezone(pytz.UTC) 
        end_time_utc = end_time.astimezone(pytz.UTC)
        
        rates = mt5.copy_rates_range(symbol, mt5.TIMEFRAME_M1, start_time_utc, end_time_utc)
        df = pd.DataFrame(rates)
        df['time'] = pd.to_datetime(df['time'], unit='s', utc=True)
        
        # Add volume column if it doesn't exist
        if 'volume' not in df.columns:
            df['volume'] = 0
            
        df = calculate_atr(df)
        return df
    else:
        # Fetch data from Polygon for other symbols
        api_symbol = symbol  
        aggs = []
        retry_attempts = 0
        current_date = start_time.astimezone(pytz.UTC).date()
        end_date = end_time.astimezone(pytz.UTC).date()

        while current_date <= end_date and retry_attempts < 10:
            try:
                for a in client.list_aggs(ticker=api_symbol, multiplier=1, timespan="minute", 
                                          from_=current_date.strftime('%Y-%m-%d'), 
                                          to=(current_date + timedelta(days=1)).strftime('%Y-%m-%d'), 
                                          limit=50000):
                    aggs.append(a)

                current_date += timedelta(days=1)
                retry_attempts = 0  
            except Exception as e:
                print(f"Error fetching data for {api_symbol} on {current_date}: {e}")
                if '429' in str(e):
                    retry_attempts += 1
                    sleep_time = 2 ** retry_attempts
                    print(f"Retrying in {sleep_time} seconds...")
                    time.sleep(sleep_time)
                else:
                    break

        if not aggs:
            print(f"No OHLC data available for {symbol}")
            return None

        df = pd.DataFrame([{
            'time': pd.to_datetime(agg.timestamp, unit='ms', utc=True),
            'open': agg.open,
            'high': agg.high,
            'low': agg.low,
            'close': agg.close,
            'volume': getattr(agg, 'volume', 0)  # Use 0 if volume doesn't exist
        } for agg in aggs])

        df = df[(df['time'] >= start_time.astimezone(pytz.UTC)) & 
                (df['time'] <= end_time.astimezone(pytz.UTC))]
        df.sort_values('time', inplace=True)
        df = calculate_atr(df)

        return df

def calculate_atr(df, period=14):
    df['tr1'] = df['high'] - df['low']
    df['tr2'] = (df['high'] - df['close'].shift()).abs()
    df['tr3'] = (df['low'] - df['close'].shift()).abs()
    df['tr'] = df[['tr1', 'tr2', 'tr3']].max(axis=1)
    df['atr'] = df['tr'].rolling(window=period).mean()
    return df

def parse_excel_input(excel_input):
    excel_input = excel_input.replace('^I', '\t')
    lines = excel_input.strip().split('\n')
    events = []
    date_time_inputs = {}
    event_data_list = {}
    event_symbols = {}

    current_event_name = ""
    current_event_time = ""
    current_event_description = ""

    for line in lines:
        if line.startswith('Event'):
            parts = line.split('\t')
            current_event_name = parts[1].strip()
            events.append({'name': current_event_name})
        elif line.startswith('Actual:'):
            parts = line.split('\t')
            symbol1 = parts[-2].strip()
            symbol2 = parts[-1].strip()
            event_symbols[current_event_name] = {"polygon": [symbol1, symbol2], "mt5": [symbol1, symbol2]}
        elif line.startswith('Forecast:'):
            parts = line.split('\t')
            symbol1 = parts[-2].strip()
            symbol2 = parts[-1].strip()
            event_symbols[current_event_name]["polygon"] = [symbol1, symbol2]
        elif line.startswith('Previous:'):
            parts = line.split('\t')
            current_event_description = parts[-1].strip()
        elif line.startswith('Time (GMT):'):
            parts = line.split('\t')
            current_event_time = parts[1].strip()
        elif re.match(r'\d{2}-\w{3}-\d{2}', line):
            parts = line.split('\t')
            date_str = parts[0].strip()
            date = datetime.strptime(date_str, '%d-%b-%y').strftime('%Y-%m-%d')
            actual = parts[1].strip().replace('%', '')
            forecast = parts[2].strip().replace('%', '') or "NA"
            previous = parts[3].strip().replace('%', '') or "NA"
            time_str = parts[4].strip()

            if date not in date_time_inputs:
                date_time_inputs[date] = {
                    "date": date,
                    "event": current_event_time + "3:00",
                }
                event_data_list[date] = []

            event_data_list[date].append({
                "name": current_event_name,
                "actual": actual,
                "forecast": forecast,
                "previous": previous,
                "time": time_str,
                "description": current_event_description
            })

    return date_time_inputs, event_data_list, events, event_symbols

# Initialize presentation
template_path = "/Users/pawan/Desktop/Quantwater Tech Investments/Research & Development/Research/Blog/Slides/template3.pptx"
prs = Presentation(template_path)
generated_image_filenames = []

def add_event_name_slide(event_name):
    slide_layout = prs.slide_layouts[2]  # Layout Index: 2, Layout Name: 1_Main Title 02
    slide = prs.slides.add_slide(slide_layout)
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "[Event Name]" in shape.text:
                shape.text = event_name

def add_event_description_slide(event_name, event_data):
    slide_layout = prs.slide_layouts[3]  # Layout Index: 3, Layout Name: 3_Main Title 02
    slide = prs.slides.add_slide(slide_layout)
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "[Event Name]" in shape.text:
                shape.text = event_name
            elif "Description" in shape.text:
                shape.text = event_data[0]["description"]
            elif "Release Date" in shape.text:
                table = shape.table
                for i, event in enumerate(event_data):
                    table.cell(i + 1, 0).text = event["time"]
                    table.cell(i + 1, 1).text = event["actual"]
                    table.cell(i + 1, 2).text = event["forecast"]
                    table.cell(i + 1, 3).text = event["previous"]

def add_asset_divider_slide(event_name, asset_name, layout_index):
    slide_layout = prs.slide_layouts[layout_index]  # Layout Index: 4 or 6
    slide = prs.slides.add_slide(slide_layout)
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "[Event Name]" in shape.text:
                shape.text = event_name
            elif "[asset name" in shape.text:
                shape.text = asset_name

def add_image_slide(image_filename, layout_index):
    slide_layout = prs.slide_layouts[layout_index]  # Layout Index: 5 or 7
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(image_filename, Inches(0), Inches(0), height=prs.slide_height)

def process_event(event, date_time_inputs, event_data_list, event_symbols):
    event_summary = {
        'volatility_units': [],
        'max_movements': [],
        'min_movements': [],
        'pre_atr_volatilities': []
    }

    # Add Event Name Slide
    add_event_name_slide(event['name'])

    # Add Event Description Slide
    event_dates = [
        date for date in sorted(date_time_inputs.keys(), key=lambda x: datetime.strptime(x, '%Y-%m-%d'))
        if any(e['name'] == event['name'] for e in event_data_list[date])
    ]
    event_data = [e for e in event_data_list[event_dates[0]] if e['name'] == event['name']]
    add_event_description_slide(event['name'], event_data)

    # Process for each symbol
    for symbol_idx, symbol in enumerate(event_symbols[event['name']]['polygon']):
        # Add Asset Divider Slide
        if symbol_idx == 0:
            add_asset_divider_slide(event['name'], symbol, 4)  # Layout Index: 4
        else:
            add_asset_divider_slide(event['name'], symbol, 6)  # Layout Index: 6

        # Process event data and generate image slides
        for idx, date_str in enumerate(event_dates):
            event_data = next((e for e in event_data_list[date_str] if e['name'] == event['name']), None)
            if not event_data:
                continue

            # Fetch data and generate image slides
            event_time_str = event_data["time"]
            event_date = datetime.strptime(date_str, '%Y-%m-%d').date()
            event_time_obj = datetime.strptime(event_time_str, '%H:%M').time()
            event_time = pytz.timezone('Etc/GMT').localize(datetime.combine(event_date, event_time_obj))

            start_time = event_time - timedelta(hours=2)
            end_time = event_time + timedelta(hours=2)

            current_spread = get_current_spread(event_symbols[event['name']]['mt5'][symbol_idx])
            if current_spread is None:
                current_spread = estimate_typical_spread(event_symbols[event['name']]['mt5'][symbol_idx].split('.')[0])

            df = fetch_data(symbol, start_time, end_time)
            if df is not None and not df.empty:
                # Convert event_time to match df timezone (UTC)
                event_time_utc = event_time.astimezone(pytz.UTC)

                # Generate OHLC chart and save image
                image_filename, volatility_unit, max_movement, min_movement, pre_atr_volatility = plot_ohlc(
                    df, symbol, event_time_utc, [event_data],
                    atr_multiple, current_spread, atr_period,
                    is_last_chart=(idx == len(event_dates) - 1),
                    start_time=start_time.astimezone(pytz.UTC),
                    end_time=end_time.astimezone(pytz.UTC)
                )

                # Add Image Slide
                if symbol_idx == 0:
                    add_image_slide(image_filename, 5)  # Layout Index: 5
                else:
                    add_image_slide(image_filename, 7)  # Layout Index: 7

                # Update event summary
                event_summary['volatility_units'].append(volatility_unit)
                event_summary['max_movements'].append(max_movement / volatility_unit)
                event_summary['min_movements'].append(min_movement / volatility_unit)
                event_summary['pre_atr_volatilities'].append(pre_atr_volatility)

                # Clean up temporary image files
                generated_image_filenames.append(image_filename)

    return event_summary

# User inputs
atr_multiple = 1.5
atr_period = 14

print("Please paste your Excel-formatted event data (press Enter twice when finished):")
excel_input = ""
while True:
    line = input()
    if line.strip() == "":
        break
    excel_input += line + "\n"

# Parse the Excel input
try:
    date_time_inputs, event_data_list, events, event_symbols = parse_excel_input(excel_input)
except ValueError as e:
    print(f"Error: {e}")
    exit(1)

# Process events
for event in events:
    event_summary = process_event(event, date_time_inputs, event_data_list, event_symbols)

output_pptx_path = "/Users/pawan/Desktop/Quantwater Tech Investments/Research & Development/Research/Blog/Slides/analyzer21.pptx"
prs.save(output_pptx_path)

# Clean up temporary image files
for image_filename in generated_image_filenames:
    os.remove(image_filename)

# Shutdown MT5
mt5.shutdown()